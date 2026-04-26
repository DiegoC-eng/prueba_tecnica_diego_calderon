"""
Generador de presentación ejecutiva en inglés (Bloque 5)
Autor: Diego Alberto Calderón Calderón
5 slides para VP de Operaciones
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colores Walmart
BLUE   = RGBColor(0x00, 0x53, 0xE2)
SPARK  = RGBColor(0xFF, 0xC2, 0x20)
RED    = RGBColor(0xEA, 0x11, 0x00)
GREEN  = RGBColor(0x2A, 0x87, 0x03)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x74, 0x76, 0x7C)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT  = RGBColor(0xE8, 0xF0, 0xFD)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Blank layout
blank_layout = prs.slide_layouts[6]  # Blank


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
            font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
            italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=14, color=DARK, bullet_char='▸ '):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet_char + b
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def make_header(slide, title, subtitle=None):
    # Blue top bar
    add_rect(slide, 0, 0, 13.33, 1.1, BLUE)
    add_text(slide, title, 0.3, 0.08, 12, 0.65,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.72, 12, 0.4,
                 font_size=13, color=SPARK, align=PP_ALIGN.LEFT)
    # Bottom bar
    add_rect(slide, 0, 7.2, 13.33, 0.3, BLUE)
    add_text(slide, 'Retail Chain Centroamerica | Confidential | April 2026',
             0.3, 7.2, 10, 0.3, font_size=9, color=WHITE)


# ============================================================
# SLIDE 1 — EXECUTIVE SUMMARY
# ============================================================
slide1 = prs.slides.add_slide(blank_layout)
make_header(slide1, 'Executive Summary', '18 Months | 40 Stores | 5 Countries | $48.7M GMV')

# 3 key finding boxes
boxes = [
    (BLUE,  'STORE PERFORMANCE',
     ['Comp stores grew GMV avg +8.3% YoY',
      'Hipermercado is the most volatile format (peak: Christmas +34%)',
      'Top 3 countries: CR, GT, SV above comp sales target']),
    (GREEN, 'LOYALTY & CUSTOMER',
     ['40.2% of transactions use loyalty card',
      'Retention M+1 ranges 30–38% across cohorts',
      'Biggest drop: M+1 to M+2 (−15 pp) — intervention needed']),
    (RED,   'RISKS IDENTIFIED',
     ['A/B test: new display shows NO significant GMV lift (p=0.24)',
      'TIENDA_012: 7-day data gap — possible POS failure',
      '231 items with price=0 outside promo — catalog error']),
]

for i, (color, title, bullets) in enumerate(boxes):
    x = 0.25 + i * 4.35
    add_rect(slide1, x, 1.25, 4.1, 0.5, color)
    add_text(slide1, title, x + 0.1, 1.3, 3.9, 0.42,
             font_size=13, bold=True, color=WHITE)
    add_rect(slide1, x, 1.75, 4.1, 4.5, LIGHT)
    add_bullet_box(slide1, bullets, x + 0.12, 1.85, 3.9, 4.2,
                   font_size=12.5, color=DARK)

# GMV stat
add_rect(slide1, 0.25, 6.3, 12.8, 0.75, SPARK)
add_text(slide1, 'Total GMV: $48.7M  |  171K transactions  |  Avg ticket: $278  |  Return rate: 2.03%',
         0.35, 6.35, 12.5, 0.6, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — STORE PERFORMANCE
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)
make_header(slide2, 'Store Performance', 'Comp Sales YoY | GMV per m² | Format Ranking')

# Left column: Comp Sales
add_rect(slide2, 0.2, 1.2, 6.1, 5.7, LIGHT)
add_text(slide2, 'Comp Sales Growth (YoY)', 0.3, 1.25, 5.9, 0.45,
         font_size=14, bold=True, color=BLUE)

comp_data = [
    ('HIPERMERCADO',  '+11.2%', GREEN),
    ('SUPERMERCADO',  '+7.8%',  GREEN),
    ('DESCUENTO',     '+5.1%',  GREEN),
    ('EXPRESS',       '+3.2%',  GRAY),
]
for i, (fmt, val, color) in enumerate(comp_data):
    y = 1.85 + i * 1.1
    add_rect(slide2, 0.3, y, 3.0, 0.7, color)
    add_text(slide2, fmt, 0.4, y + 0.05, 2.8, 0.6, font_size=12, bold=True, color=WHITE)
    add_text(slide2, val, 3.4, y + 0.05, 2.5, 0.6,
             font_size=22, bold=True, color=color)

# Right column: GMV/m2
add_rect(slide2, 6.6, 1.2, 6.5, 5.7, LIGHT)
add_text(slide2, 'GMV per m² — Last Quarter (Apr–Jun 2025)',
         6.7, 1.25, 6.3, 0.45, font_size=14, bold=True, color=BLUE)

bullets_perf = [
    '8 stores flagged as BAJO_RENDIMIENTO (P25 of their format)',
    'EXPRESS stores show highest GMV/m² efficiency despite low absolute GMV',
    'Bottom 3 HIPERMERCADO stores: 42% below format median',
    'Recommended action: store visit + surtido review for underperformers',
    'Productivity gap between top and bottom quartile: $18.4/m²',
]
add_bullet_box(slide2, bullets_perf, 6.7, 1.85, 6.2, 4.8,
               font_size=12.5, color=DARK)


# ============================================================
# SLIDE 3 — OPPORTUNITIES
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
make_header(slide3, 'Opportunities', 'Underperforming Stores | Low GMROI Vendors | Loyalty')

opp_items = [
    (BLUE,  '1. Underperforming Stores',
     '8 stores below P25 GMV/m² in their format. Potential GMV uplift '
     'if brought to median: ~$1.2M/quarter. Priority: 3 HIPERMERCADO stores in HN & NI.'),
    (SPARK, '2. Low GMROI Vendors',
     'Vendors with GMROI < 1.0 are generating LESS gross margin than their cost. '
     'Renegotiate terms or delist products. Focus on Tier C vendors in '
     'Electronics & Apparel categories.'),
    (GREEN, '3. Loyalty Retention Gap',
     'M+1 to M+2 retention drops 15 pp on average. '
     'A targeted coupon campaign at week 3 post-second-purchase could recover '
     '5–10 pp retention → ~$180K incremental annual GMV from loyalty segment.'),
    (GRAY,  '4. Digital Payment Adoption',
     'HN & NI average 40% cash transactions vs 25% in CR. '
     'Digital payment adoption reduces operational cost by ~$0.8/transaction. '
     'Partner with mobile wallets (Tigo Money) in low-banked markets.'),
]

for i, (color, title, body) in enumerate(opp_items):
    y = 1.2 + i * 1.4
    add_rect(slide3, 0.2, y, 0.15, 1.1, color)
    add_text(slide3, title, 0.45, y, 4.0, 0.45,
             font_size=13, bold=True, color=color)
    txb = slide3.shapes.add_textbox(
        Inches(0.45), Inches(y + 0.45), Inches(12.6), Inches(0.85))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(12)
    run.font.color.rgb = DARK


# ============================================================
# SLIDE 4 — RISKS
# ============================================================
slide4 = prs.slides.add_slide(blank_layout)
make_header(slide4, 'Risks', 'Stock-Outs | Retention Drop | Data Quality')

risks = [
    ('Stock-Out Impact',
     ['Electronics & Toys show highest estimated GMV at risk from low-rotation SKUs',
      'Gap analysis identified items with 3+ consecutive days without sales in active stores',
      'Estimated annual GMV at risk: to be quantified with full OOS detection pipeline',
      'Seasonality peak (Dec) amplifies risk: surtido must be ready by Nov 1st']),
    ('Retention Drop Risk',
     ['Recent cohorts (Oct–Nov 2024) show lower M+1 retention than older cohorts',
      'If trend continues, loyalty program ROI declines without intervention',
      'M+6 retention: only ~18–22% of original cohort remains active',
      'Risk: loyalty card penetration stalls at 40% without new value proposition']),
    ('Data Quality Risks',
     ['50 transactions occurred BEFORE store opening_date — POS configuration error',
      '1 vendor_id orphaned in product catalog — GMROI analysis incomplete',
      '1,745 transactions with total_amount discrepancy vs item sum (max gap: $202)',
      'TIENDA_008 & TIENDA_037: dual A/B assignment — random assignment process must be fixed']),
]

for i, (title, bullets) in enumerate(risks):
    x = 0.25 + i * 4.35
    add_rect(slide4, x, 1.2, 4.1, 0.5, RED if i == 0 else (SPARK if i == 1 else GRAY))
    color = RED if i == 0 else (SPARK if i == 1 else GRAY)
    add_text(slide4, title, x + 0.1, 1.25, 3.9, 0.42,
             font_size=12, bold=True, color=WHITE)
    add_rect(slide4, x, 1.7, 4.1, 5.3, LIGHT)
    tc = DARK if i == 2 else DARK
    add_bullet_box(slide4, bullets, x + 0.12, 1.8, 3.9, 5.0,
                   font_size=11.5, color=tc)


# ============================================================
# SLIDE 5 — RECOMMENDATIONS
# ============================================================
slide5 = prs.slides.add_slide(blank_layout)
make_header(slide5, 'Recommendations', 'What to do | Who | By When')

recs = [
    (BLUE,  'Fix A/B Test Process',
     'Owner: Analytics Team',
     'By: May 15, 2026',
     'Remove TIENDA_008 & TIENDA_037 from active experiments. '
     'Implement randomization protocol with pre-test balance check '
     'before any future experiment launches.'),
    (GREEN, 'Launch Loyalty Intervention at M+2',
     'Owner: CRM / Marketing',
     'By: June 1, 2026',
     'Deploy personalized coupon at day 21 post-second-purchase '
     'for all loyalty customers. Target: recover 5 pp retention. '
     'Expected GMV impact: +$180K/year.'),
    (RED,   'Store Productivity Action Plan',
     'Owner: Regional Operations VP',
     'By: May 30, 2026',
     '8 BAJO_RENDIMIENTO stores need store visit in 30 days. '
     'Focus: surtido optimization, staff training, local promo calendar. '
     'Target: bring all stores to format median GMV/m² by Q3.'),
    (SPARK, 'Renegotiate Low-GMROI Vendors',
     'Owner: Buying / Procurement',
     'By: Next Vendor Review Cycle (Q3)',
     'Identify all Tier C vendors with GMROI < 1.0. '
     'Negotiate cost reduction or delist. Reallocate shelf space '
     'to high-velocity, high-margin SKUs.'),
    (GRAY,  'Implement Daily Store Data Monitor',
     'Owner: Data Engineering',
     'By: June 15, 2026',
     'Automated alert when any store has 0 transactions for >24h. '
     'Fixes root cause of TIENDA_012 gap and pre-opening data issue. '
     'Alert via Slack + dashboard red indicator.'),
]

for i, (color, title, owner, deadline, body) in enumerate(recs):
    y = 1.2 + i * 1.18
    add_rect(slide5, 0.2, y, 0.12, 0.95, color)
    add_text(slide5, title, 0.42, y, 5.5, 0.4,
             font_size=12, bold=True, color=color)
    add_text(slide5, f'{owner}  |  {deadline}',
             0.42, y + 0.38, 5.5, 0.35, font_size=10, color=GRAY, italic=True)
    txb = slide5.shapes.add_textbox(
        Inches(6.2), Inches(y), Inches(6.9), Inches(0.95))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(11)
    run.font.color.rgb = DARK


# Save
prs.save('bloque5_presentacion_EN.pptx')
print('PPTX guardado: bloque5_presentacion_EN.pptx')

# Convert to PDF using matplotlib (each slide as PNG, then PDF)
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation as PRS

# We saved the PPTX. Note: for the PDF version, we'll convert via alternative
# For now, rename and note that LibreOffice or PowerPoint can convert
print('Nota: Para generar bloque5_presentacion_EN.pdf, abrir el .pptx en PowerPoint > Guardar como PDF')
print('O usar: libreoffice --headless --convert-to pdf bloque5_presentacion_EN.pptx')

import shutil
shutil.copy('bloque5_presentacion_EN.pptx', 'bloque5_presentacion_EN_backup.pptx')
print('Done!')
